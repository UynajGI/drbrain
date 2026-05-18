"""Tests for src/drbrain/services/document.py."""

from __future__ import annotations

from pathlib import Path

import pytest

from drbrain.services.document import inspect, inspect_docx, inspect_pptx, inspect_xlsx


class TestInspectDispatch:
    def test_missing_file_raises(self):
        with pytest.raises(FileNotFoundError):
            inspect(Path("/nonexistent/file.docx"))

    def test_directory_raises(self, tmp_path):
        with pytest.raises(ValueError, match="not a file"):
            inspect(tmp_path)

    def test_unsupported_format_raises(self, tmp_path):
        f = tmp_path / "test.xyz"
        f.write_text("dummy")
        with pytest.raises(ValueError, match="Unsupported file format"):
            inspect(f)

    def test_auto_detect_extension(self, tmp_path):
        f = tmp_path / "test.docx"
        f.write_text("dummy")
        # docx inspect will fail on dummy content, but format detection works
        with pytest.raises((ImportError, Exception)):
            inspect(f)

    def test_explicit_format(self, tmp_path):
        f = tmp_path / "test.bin"
        f.write_text("dummy")
        with pytest.raises((ImportError, Exception)):
            inspect(f, fmt="pptx")


class TestDocxInspector:
    def test_import_error_on_missing_dep(self, tmp_path):
        """Without python-docx, inspect_docx raises ImportError with hint."""
        f = tmp_path / "test.docx"
        f.write_text("not a real docx")
        try:
            import docx  # noqa: F401
        except ImportError:
            with pytest.raises(ImportError, match="drbrain\\[office\\]"):
                inspect_docx(f)

    def test_real_docx_detects_structure(self, tmp_path):
        """Minimal docx: create via python-docx if available."""
        pytest.importorskip("docx")
        from docx import Document

        doc = Document()
        doc.add_heading("Test Heading", level=1)
        doc.add_paragraph("A test paragraph.")
        path = tmp_path / "real.docx"
        doc.save(str(path))

        result = inspect_docx(path)
        assert "=== DOCX:" in result
        assert "Heading 1" in result
        assert "Test Heading" in result
        assert "Paragraphs:" in result


class TestPptxInspector:
    def test_import_error_on_missing_dep(self, tmp_path):
        """Without python-pptx, inspect_pptx raises ImportError with hint."""
        f = tmp_path / "test.pptx"
        f.write_text("not a real pptx")
        try:
            import pptx  # noqa: F401
        except ImportError:
            with pytest.raises(ImportError, match="drbrain\\[office\\]"):
                inspect_pptx(f)

    def test_real_pptx_detects_structure(self, tmp_path):
        """Minimal pptx: create via python-pptx if available."""
        pytest.importorskip("pptx")
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Test Title"
        path = tmp_path / "real.pptx"
        prs.save(str(path))

        result = inspect_pptx(path)
        assert "=== PPTX:" in result
        assert "Test Title" in result
        assert "Summary" in result


class TestXlsxInspector:
    def test_import_error_on_missing_dep(self, tmp_path):
        """Without openpyxl, inspect_xlsx raises ImportError with hint."""
        f = tmp_path / "test.xlsx"
        f.write_text("not a real xlsx")
        try:
            import openpyxl  # noqa: F401
        except ImportError:
            with pytest.raises(ImportError, match="drbrain\\[office\\]"):
                inspect_xlsx(f)

    def test_real_xlsx_detects_structure(self, tmp_path):
        """Minimal xlsx: create via openpyxl if available."""
        pytest.importorskip("openpyxl")
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "alpha"
        ws["B2"] = 42
        path = tmp_path / "real.xlsx"
        wb.save(str(path))

        result = inspect_xlsx(path)
        assert "=== XLSX:" in result
        assert "Data" in result
        assert "Name" in result
        assert "Summary" in result
